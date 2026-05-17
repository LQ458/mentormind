"""
MentorMind Presentation — Light Professional Edition
Clean white backgrounds, blue accents, embedded diagrams, 10 slides.
"""

import os
DIAGRAM_DIR = os.path.join(os.path.dirname(__file__), "diagrams")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Light theme palette ──
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
TEXT    = RGBColor(0x0F, 0x17, 0x2A)
BODY    = RGBColor(0x33, 0x40, 0x55)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
BLUE    = RGBColor(0x25, 0x63, 0xEB)
BLUE_L  = RGBColor(0xDB, 0xEA, 0xFE)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
GREEN   = RGBColor(0x05, 0x96, 0x69)
ORANGE  = RGBColor(0xEA, 0x58, 0x0C)
BORDER  = RGBColor(0xE2, 0xE8, 0xF0)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Helpers ──
def bg(s, c=BG):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def bar(s, c=BLUE):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Pt(4))
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()

def line(s, l, t, w, c=BLUE, h=Pt(3)):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()

def txt(s, l, t, w, h, text, sz=Pt(32), clr=TEXT, bold=True, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = sz; p.font.color.rgb = clr
    p.font.bold = bold; p.alignment = align
    return tf

def bullets(s, l, t, w, h, items, sz=Pt(16), clr=BODY):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = sz; p.font.color.rgb = clr
        p.space_before = Pt(16) if i > 0 else Pt(0)
    return tf

def note(s, text):
    s.notes_slide.notes_text_frame.text = text

def pgnum(s, n):
    tb = s.shapes.add_textbox(W - Inches(1), H - Inches(0.5), Inches(0.8), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = f"{n} / 10"; p.font.size = Pt(9); p.font.color.rgb = MUTED; p.alignment = PP_ALIGN.RIGHT

def embed(s, name, l, t, w, h):
    path = os.path.join(DIAGRAM_DIR, f"{name}.png")
    if os.path.exists(path):
        s.shapes.add_picture(path, l, t, w, h)

def card(s, l, t, w, h, number, label, color=BLUE):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = CARD_BG
    c.line.color.rgb = BORDER; c.line.width = Pt(1)
    # shadow effect via a slightly offset dark rect behind
    # number
    txt(s, l + Inches(0.2), t + Inches(0.2), w - Inches(0.4), h * 0.55,
        number, Pt(26), color, True, PP_ALIGN.CENTER)
    txt(s, l + Inches(0.2), t + h * 0.55, w - Inches(0.4), h * 0.4,
        label, Pt(11), BODY, False, PP_ALIGN.CENTER)

def circled_num(s, l, t, d, n, c=BLUE):
    ci = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, d, d)
    ci.fill.solid(); ci.fill.fore_color.rgb = c; ci.line.fill.background()
    p = ci.text_frame.paragraphs[0]
    p.text = n; p.font.size = Pt(16); p.font.color.rgb = WHITE
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER

# ── SLIDE 1: Title ──
s = prs.slides.add_slide(blank)
bg(s, WHITE)
# Blue accent stripe left
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), H)
r.fill.solid(); r.fill.fore_color.rgb = BLUE; r.line.fill.background()

txt(s, Inches(1.5), Inches(1.6), Inches(10), Inches(1.5), "MentorMind", Pt(52), TEXT, True, PP_ALIGN.LEFT)
txt(s, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8), "AI Powered Personalized Study Plans", Pt(24), BLUE, False, PP_ALIGN.LEFT)
line(s, Inches(1.5), Inches(3.8), Inches(2), BLUE, Pt(3))
txt(s, Inches(1.5), Inches(4.3), Inches(10), Inches(0.6), "Senior Project Presentation    Leo Qin", Pt(16), MUTED, False, PP_ALIGN.LEFT)
note(s, "Hi everyone. MentorMind is an AI platform that builds personalized study plans for Chinese students. In 10 minutes I will show you how we diagnose knowledge gaps and generate adaptive learning paths, plus a quick live demo.")

# ── SLIDE 2: Problem ──
s = prs.slides.add_slide(blank)
bg(s); bar(s)
txt(s, Inches(1), Inches(0.6), Inches(11.3), Inches(1.1), "The Problem", Pt(34), TEXT)
line(s, Inches(1), Inches(1.4), Inches(1.8), BLUE, Pt(3))

card(s, Inches(1), Inches(2.2), Inches(3.5), Inches(1.5), "46M+", "Students in exam track education", BLUE)
card(s, Inches(4.9), Inches(2.2), Inches(3.5), Inches(1.5), "$2K to 15K", "Annual tutoring cost per family", ORANGE)
card(s, Inches(8.8), Inches(2.2), Inches(3.5), Inches(1.5), "50 to 1", "Student to teacher ratio", PURPLE)

bullets(s, Inches(1), Inches(4.3), Inches(11.3), Inches(2.5), [
    "Private tutoring is expensive but remains one size fits all",
    "Teachers cannot diagnose individual gaps for 50 students",
    "Students waste hours reviewing what they already know",
])
note(s, "China has 46 million exam track students. Private tutoring costs families thousands per year but remains generic. A teacher with 50 students cannot deeply diagnose each one. Students end up restudying mastered topics while their real gaps go unaddressed.")
pgnum(s, 2)

# ── SLIDE 3: Solution ──
s = prs.slides.add_slide(blank)
bg(s); bar(s)
txt(s, Inches(1), Inches(0.5), Inches(11.3), Inches(0.9), "How MentorMind works", Pt(34), TEXT)
line(s, Inches(1), Inches(1.1), Inches(1.8), BLUE, Pt(3))

steps = [
    ("1", "Conversational diagnostic", "AI chats with the student to uncover real knowledge gaps"),
    ("2", "Knowledge graph", "Builds a per student map of concepts and their dependencies"),
    ("3", "Manager Critic AI", "Generates study plans then evaluates quality before delivery"),
    ("4", "Active learning modes", "Seminar debates, simulations, oral defense, memory challenges"),
]
for i, (n, title, desc) in enumerate(steps):
    y = Inches(1.7 + i * 1.2)
    circled_num(s, Inches(1.2), y + Pt(4), Inches(0.45), n, BLUE)
    txt(s, Inches(2.1), y, Inches(10), Inches(0.45), title, Pt(18), TEXT)
    txt(s, Inches(2.1), y + Pt(28), Inches(10), Inches(0.35), desc, Pt(13), BODY, False)

note(s, "Four steps. First an AI chats with the student, adapting questions based on answers. Second we build a personal knowledge graph. Third a manager critic AI pair generates and checks study plans. Fourth the student engages through five active learning modes instead of passive reading.")
pgnum(s, 3)

# ── SLIDE 4: Architecture ──
s = prs.slides.add_slide(blank)
bg(s); bar(s)
txt(s, Inches(1), Inches(0.4), Inches(11.3), Inches(0.8), "System Architecture", Pt(32), TEXT)
line(s, Inches(1), Inches(0.9), Inches(1.5), BLUE, Pt(3))

bullets(s, Inches(1), Inches(1.4), Inches(5.5), Inches(5.2), [
    "Next.js 14 frontend with Clerk authentication",
    "FastAPI v2 backend serving ~90 endpoints",
    "Three isolated Celery worker queues",
    "PostgreSQL for all data, Redis for caching",
    "SiliconFlow API with DeepSeek R1 and V3",
    "FunASR for speech, PaddleOCR for image text",
], Pt(15))
embed(s, "architecture", Inches(7), Inches(1.2), Inches(5.6), Inches(5.3))
note(s, "Next.js frontend with Clerk auth. FastAPI backend with about 90 endpoints. Three Celery queues keep heavy tasks isolated. PostgreSQL stores everything including the knowledge graph. Redis handles caching and job orchestration. AI through SiliconFlow at under one cent per thousand tokens. FunASR and PaddleOCR handle Chinese speech and images.")
pgnum(s, 4)

# ── SLIDE 5: Knowledge Graph ──
s = prs.slides.add_slide(blank)
bg(s); bar(s)
txt(s, Inches(1), Inches(0.4), Inches(11.3), Inches(0.8), "Personal Knowledge Graph", Pt(32), TEXT)
line(s, Inches(1), Inches(0.9), Inches(1.5), BLUE, Pt(3))

bullets(s, Inches(1), Inches(1.4), Inches(5.5), Inches(5.2), [
    "Audio uploads processed by FunASR",
    "Image uploads processed by PaddleOCR",
    "LLM extracts concepts and relationships",
    "Graph stored per user in PostgreSQL",
    "D3.js interactive force directed layout",
    "Mastery colors: green, yellow, red",
], Pt(15))
embed(s, "knowledge_graph", Inches(7), Inches(1.2), Inches(5.6), Inches(5.3))
note(s, "Students upload audio or images. FunASR and PaddleOCR extract text. An LLM identifies concepts and their relationships such as prerequisites and containment. Everything is stored in a per user knowledge graph in PostgreSQL. The D3.js frontend colors nodes by mastery level. Red nodes become the study plan priority.")
pgnum(s, 5)

# ── SLIDE 6: Manager Critic ──
s = prs.slides.add_slide(blank)
bg(s); bar(s)
txt(s, Inches(1), Inches(0.4), Inches(11.3), Inches(0.8), "Manager Critic AI Planning Loop", Pt(32), TEXT)
line(s, Inches(1), Inches(0.9), Inches(1.5), BLUE, Pt(3))

bullets(s, Inches(1), Inches(1.4), Inches(5.5), Inches(5.2), [
    "DeepSeek R1 Manager generates the draft plan",
    "DeepSeek V3 Critic evaluates on 5 dimensions",
    "Clarity, accuracy, pedagogy, engagement, difficulty",
    "Scored 0 to 1 with quality threshold at 0.8",
    "Below threshold triggers regeneration up to 3 times",
    "Four subagents produce quizzes and flashcards",
], Pt(15))
embed(s, "manager_critic", Inches(7), Inches(1.2), Inches(5.6), Inches(5.3))
note(s, "R1 acts as manager generating the study plan with units and objectives. V3 acts as critic evaluating on five dimensions including clarity, accuracy, pedagogical effectiveness, engagement, and difficulty. Threshold is 0.8. Below that the critic sends feedback and R1 regenerates up to three times. Plans that pass go to the student. Four subagents produce the actual content.")
pgnum(s, 6)

# ── SLIDE 7: Process First ──
s = prs.slides.add_slide(blank)
bg(s); bar(s)
txt(s, Inches(1), Inches(0.4), Inches(11.3), Inches(0.8), "Learning by doing, not reading", Pt(32), TEXT)
line(s, Inches(1), Inches(0.9), Inches(1.5), BLUE, Pt(3))

bullets(s, Inches(1), Inches(1.4), Inches(5.5), Inches(5.2), [
    "Seminar: three AI roles debate, you evaluate",
    "Simulation: decision scenarios with real concepts",
    "Oral Defense: expert panel questions your reasoning",
    "Memory Challenge: timed retrieval sprints",
    "Error Audit: identify flaws in wrong answers",
    "Backed by forgetting curve spaced repetition",
], Pt(15))
embed(s, "process_flow", Inches(7), Inches(1.2), Inches(5.6), Inches(5.3))
note(s, "What sets MentorMind apart. Seminar has three AI personas debate and the student evaluates. Simulation presents real scenarios. Oral Defense fires questions from a panel of experts. Memory Challenge is a timed sprint. Error Audit finds flaws in a wrong answer. All backed by a forgetting curve scheduler with proactive review notifications.")
pgnum(s, 7)

# ── SLIDE 8: Pivot ──
s = prs.slides.add_slide(blank)
bg(s); bar(s)
txt(s, Inches(1), Inches(0.5), Inches(11.3), Inches(0.9), "The Pivot", Pt(34), TEXT)
line(s, Inches(1), Inches(1.1), Inches(1.5), BLUE, Pt(3))

# Two column layout
bullets(s, Inches(1), Inches(1.8), Inches(5.5), Inches(4.5), [
    "Weeks 1 to 3: built a 5 stage AI video pipeline",
    "Used Manim for math animations like 3Blue1Brown",
    "CJK font rendering was inconsistent in production",
    "User research showed students wanted personalization",
], Pt(15))

bullets(s, Inches(7), Inches(1.8), Inches(5.3), Inches(4.5), [
    "Week 4: pivoted to study plan generation",
    "Hid video UI, kept pipeline code intact",
    "Built circuit breaker with multi provider fallback",
    "Pivot unlocked the real user need for personalization",
], Pt(15))

# Arrow between columns
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.8), Inches(3.2), Inches(1.2), Inches(0.5))
arrow.fill.solid(); arrow.fill.fore_color.rgb = ORANGE; arrow.line.fill.background()
txt(s, Inches(5.8), Inches(3.7), Inches(1.2), Inches(0.3), "PIVOT", Pt(9), WHITE, True, PP_ALIGN.CENTER)

note(s, "We built a full video generation pipeline with Manim animations. It worked but CJK font rendering was inconsistent. User research showed students wanted personalization more than polished videos. We pivoted in week 4, hid the video UI, kept the pipeline intact, and redirected to study plans. We also built a circuit breaker with multi provider fallback for API resilience.")
pgnum(s, 8)

# ── SLIDE 9: Results ──
s = prs.slides.add_slide(blank)
bg(s); bar(s)
txt(s, Inches(1), Inches(0.5), Inches(11.3), Inches(0.9), "What We Shipped", Pt(34), TEXT)
line(s, Inches(1), Inches(1.1), Inches(1.5), BLUE, Pt(3))

card(s, Inches(1), Inches(1.8), Inches(5.5), Inches(1.3), "~90", "API endpoints across the full platform", BLUE)
card(s, Inches(7), Inches(1.8), Inches(5.3), Inches(1.3), "4", "Languages: Chinese  English  Japanese  Korean", PURPLE)
card(s, Inches(1), Inches(3.5), Inches(5.5), Inches(1.3), "$160 per month", "Operating budget at under one cent per K tokens", GREEN)
card(s, Inches(7), Inches(3.5), Inches(5.3), Inches(1.3), "90 days", "Telemetry retention with daily proficiency rollups", ORANGE)

bullets(s, Inches(1), Inches(5.4), Inches(11.3), Inches(1.5), [
    "Adaptive difficulty adjusts from a sliding window of three quiz scores",
    "Gaokao exam prep system built and ready for our target market",
], Pt(15))
note(s, "We shipped API v2 with about 90 endpoints across lessons, study plans, knowledge graphs, Gaokao prep, board lessons, billing, and analytics. Four languages. Operating cost around 160 dollars per month. Adaptive difficulty adjusts from three quiz scores. 90 day telemetry. And a Gaokao prep system directly relevant to our Chinese market.")
pgnum(s, 9)

# ── SLIDE 10: Future ──
s = prs.slides.add_slide(blank)
bg(s, WHITE)
# Blue accent left
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), H)
r.fill.solid(); r.fill.fore_color.rgb = BLUE; r.line.fill.background()

txt(s, Inches(1.5), Inches(1.2), Inches(10), Inches(1.0), "Next Steps", Pt(36), TEXT, True, PP_ALIGN.LEFT)
line(s, Inches(1.5), Inches(2.0), Inches(1.5), BLUE, Pt(3))

bullets(s, Inches(1.5), Inches(2.6), Inches(10), Inches(2.5), [
    "Restore video generation once CJK font rendering is solved",
    "Build a mobile app with offline study plan access",
    "Integrate with school LMS platforms via LTI 1.3",
], Pt(17), TEXT)

txt(s, Inches(1.5), Inches(5.5), Inches(10), Inches(0.7), "Thank you. Questions welcome.", Pt(22), BLUE, False, PP_ALIGN.LEFT)
note(s, "Three priorities: restore video generation once CJK fonts are solved, build a React Native mobile app for offline access, and integrate with school LMS platforms. Thank you. I am happy to take questions.")
pgnum(s, 10)

# ── Save ──
out = os.path.join(os.path.dirname(__file__), "MentorMind_Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
print("10 slides, light theme, diagrams embedded, speaker notes ready")
